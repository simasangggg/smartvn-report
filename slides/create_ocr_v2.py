#!/usr/bin/env python3
"""
OCR Presentation → PPTX (21 slides, với hình minh họa)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
PRIMARY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0x2E, 0x86, 0xDE)
ACCENT_LIGHT = RGBColor(0xD6, 0xEA, 0xF8)
SUCCESS = RGBColor(0x27, 0xAE, 0x60)
WARNING = RGBColor(0xF3, 0x9C, 0x12)
DANGER = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
GRAY_TEXT = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)

IMG = "/home/soang/.openclaw/workspace/smartvn-report/slides/ocr-images"

def img(name):
    return os.path.join(IMG, name)

def add_bg(s, color=PRIMARY):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = color

def add_bar(s):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.06))
    sh.fill.solid(); sh.fill.fore_color.rgb = ACCENT; sh.line.fill.background()

def add_footer(s):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), Inches(13.333), Inches(0.4))
    sh.fill.solid(); sh.fill.fore_color.rgb = PRIMARY; sh.line.fill.background()
    tf = sh.text_frame; tf.paragraphs[0].text = "Trích xuất Hóa đơn Tiếng Việt — PaddleOCR • VietOCR • LayoutLMv3"
    tf.paragraphs[0].font.size = Pt(10); tf.paragraphs[0].font.color.rgb = WHITE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def heading(s, text, top=Inches(0.15)):
    tb = s.shapes.add_textbox(Inches(0.5), top, Inches(12), Inches(0.8))
    tf = tb.text_frame; p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = PRIMARY

def bullets(s, items, left=Inches(0.5), top=Inches(1.1), width=Inches(12), height=Inches(5.5)):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if item.startswith("##"):
            p.text = item[2:].strip(); p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = ACCENT; p.space_before = Pt(14)
        elif item.startswith("- "):
            p.text = "• " + item[2:]; p.font.size = Pt(15); p.font.color.rgb = DARK_TEXT; p.space_before = Pt(5)
        else:
            p.text = item; p.font.size = Pt(15); p.font.color.rgb = DARK_TEXT; p.space_before = Pt(5)

def add_img(s, path, left, top, width, height=None):
    if os.path.exists(path):
        try:
            if height: s.shapes.add_picture(path, left, top, width, height)
            else: s.shapes.add_picture(path, left, top, width)
        except: pass

def box(s, l, t, w, h, fill, border, title, lines):
    tb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, Inches(0.4))
    tb.fill.solid(); tb.fill.fore_color.rgb = border; tb.line.fill.background()
    tf = tb.text_frame; tf.paragraphs[0].text = title; tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = WHITE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    cb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t+Inches(0.4), w, h-Inches(0.4))
    cb.fill.solid(); cb.fill.fore_color.rgb = fill; cb.line.color.rgb = border; cb.line.width = Pt(1)
    tf2 = cb.text_frame; tf2.word_wrap = True; tf2.margin_left = Pt(10); tf2.margin_top = Pt(5)
    for i, line in enumerate(lines):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = "• " + line; p.font.size = Pt(12); p.font.color.rgb = DARK_TEXT; p.space_before = Pt(3)

def slide(prs, title=None, bg=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if bg: add_bg(s, bg)
    else: add_bar(s); add_footer(s); heading(s, title)
    return s

def section(prs, num, title):
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, ACCENT)
    tb = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(2), Inches(1.5))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(72); p.font.bold = True; p.font.color.rgb = WHITE
    tb2 = s.shapes.add_textbox(Inches(1), Inches(4), Inches(10), Inches(1.5))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.text = title; p2.font.size = Pt(32); p2.font.bold = True; p2.font.color.rgb = WHITE
    return s

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── 1. Cover ───────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, PRIMARY)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(5), Inches(5))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.fill.fore_color.brightness = 0.3; c.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "TRÍCH XUẤT THÔNG TIN\nHÓA ĐƠN TIẾNG VIỆT TỪ ẢNH"; p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.text = "Pipeline PaddleOCR — VietOCR — LayoutLMv3"; p2.font.size = Pt(18); p2.font.color.rgb = ACCENT_LIGHT; p2.space_before = Pt(14)
    tb2 = s.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    tf2 = tb2.text_frame; tf2.paragraphs[0].text = "Học viện Công nghệ Bưu chính Viễn thông\nKhoa Công nghệ Thông tin 2 — Tháng 6, 2026"; tf2.paragraphs[0].font.size = Pt(14); tf2.paragraphs[0].font.color.rgb = ACCENT_LIGHT

    # ── 2. Outline ─────────────────────────────────────────
    s = slide(prs, "NỘI DUNG TRÌNH BÀY")
    bullets(s, ["## 1. Mở đầu — Bối cảnh bài toán", "## 2. Tổng quan — Mô hình hóa & Pipeline",
        "## 3. Phân loại & Tiến hóa — 3 thế hệ thuật toán", "## 4. Trực giác — Con người đọc hóa đơn ra sao?",
        "## 5. Toán học hóa — Công thức cốt lõi", "## 6. Khoảng cách — Lý thuyết vs Thực tế",
        "## 7. Thách thức — Rào cản & Hướng phát triển", "## 8. Kết luận"])

    # ── 3. Bối cảnh ────────────────────────────────────────

    s = slide(prs, "BỐI CẢNH BÀI TOÁN")
    bullets(s, ["- Hàng triệu giao dịch thương mại tại Việt Nam tạo ra hàng triệu tờ hóa đơn giấy",
        "- Kế toán viên gõ tay từng con số → sai sót nhập liệu thường xuyên",
        "", "## Thách thức đọc hóa đơn thực tế:",
        "- Font chữ tùy ý, góc chụp nghiêng", "- Ánh sáng không đều", "- Tiếng Việt với đầy đủ dấu thanh",
        "", "- Cần biến **ảnh hóa đơn** thành **dữ liệu có cấu trúc**"])

    # ── 4. Tình huống ──────────────────────────────────────
    s = slide(prs, "CÁC TÌNH HUỐNG THỰC TẾ")
    box(s, Inches(0.5), Inches(1.3), Inches(12), Inches(1.3), ACCENT_LIGHT, ACCENT,
        "Kịch bản A — Kế toán doanh nghiệp nhỏ", ["300 hóa đơn công tác phí, mất 2–3 ngày nhập thủ công"])
    box(s, Inches(0.5), Inches(2.8), Inches(12), Inches(1.3), RGBColor(0xE8,0xF5,0xE9), SUCCESS,
        "Kịch bản B — Ứng dụng chi tiêu cá nhân", ["Chụp ảnh hóa đơn → tự động điền ngày, siêu thị, tổng tiền"])
    box(s, Inches(0.5), Inches(4.3), Inches(12), Inches(1.3), RGBColor(0xFF,0xF3,0xE0), WARNING,
        "Kịch bản C — Kiểm toán & Tuân thủ thuế", ["Xác minh hàng nghìn hóa đơn VAT tự động"])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12), Inches(0.5))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.text = "→ Nhu cầu chung: Document Information Extraction (DIE)"; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = ACCENT; p.alignment = PP_ALIGN.CENTER

    # ── 5. Ba tầng ─────────────────────────────────────────

    s = slide(prs, "MÔ HÌNH HÓA — BA TẦNG XỬ LÝ")
    add_img(s, img("three-tiers.png"), Inches(2), Inches(1.3), Inches(9), Inches(5.5))

    # ── 6. Input/Output ────────────────────────────────────
    s = slide(prs, "ĐẦU VÀO & ĐẦU RA")
    box(s, Inches(0.5), Inches(1.3), Inches(12), Inches(1.2), ACCENT_LIGHT, ACCENT,
        "Đầu vào", ["Ảnh JPEG/PNG hóa đơn chụp bằng điện thoại (nghiêng, nhòe, ánh sáng không đều)"])
    box(s, Inches(0.5), Inches(3), Inches(12), Inches(3.5), RGBColor(0xE8,0xF5,0xE9), SUCCESS,
        "Đầu ra — JSON có ngữ nghĩa",
        ['{ "store_name": "Circle K", "date": "2024-05-27", "total": 87000',
         '  "items": [ { "name": "Cafe sua da", "qty": 2, "price": 29000 }, ... ] }'])

    # ── 7. Pipeline ────────────────────────────────────────
    s = slide(prs, "PIPELINE XỬ LÝ")
    add_img(s, img("pipeline.png"), Inches(1), Inches(1.5), Inches(11), Inches(5))

    # ── 8. Text Detection ──────────────────────────────────

    s = slide(prs, "TEXT DETECTION — Từ sliding window đến anchor-free")
    box(s, Inches(0.3), Inches(1.2), Inches(4), Inches(2.5), LIGHT_BG, GRAY_TEXT,
        "Thế hệ 1 (trước 2015)", ["MSER, SWT", "Phân tích pixel theo màu sắc", "Chỉ hoạt động với ảnh sạch"])
    box(s, Inches(0.3), Inches(3.9), Inches(4), Inches(2.5), RGBColor(0xFF,0xF3,0xE0), WARNING,
        "Thế hệ 2 (2015–2017)", ["CTPN, TextBoxes++, SegLink", "CNN-based với Anchor Boxes", "Khó phát hiện văn bản cong"])
    box(s, Inches(4.6), Inches(1.2), Inches(8.4), Inches(5.2), ACCENT_LIGHT, ACCENT,
        "★ Thế hệ 3 (2018–nay) — DB", ["DB (Differentiable Binarization)", "Xương sống của PaddleOCR",
        "Dự đoán probability map per pixel", "Ngưỡng có thể học được (differentiable threshold)", "Tối ưu cho CPU/mobile"])

    # ── 9. Text Recognition ────────────────────────────────
    s = slide(prs, "TEXT RECOGNITION — Từ template matching đến Attention")
    box(s, Inches(0.3), Inches(1.2), Inches(4), Inches(2.5), LIGHT_BG, GRAY_TEXT,
        "Thế hệ 1 (trước 2014)", ["Tesseract v3, template matching", "Tách từng ký tự", "Yêu cầu font rõ ràng"])
    box(s, Inches(0.3), Inches(3.9), Inches(4), Inches(2.5), RGBColor(0xFF,0xF3,0xE0), WARNING,
        "Thế hệ 2 (2015–2018)", ["CRNN + CTC", "CNN + BiLSTM + CTC loss", "Nền tảng OCR hiện đại"])
    box(s, Inches(4.6), Inches(1.2), Inches(8.4), Inches(5.2), ACCENT_LIGHT, ACCENT,
        "★ Thế hệ 3 (2018–nay) — VietOCR", ["Seq2Seq + Attention", "VGG/ResNet encoder + Transformer decoder",
        "Xử lý dấu thanh tiếng Việt tốt hơn", "Pretrain trên 5M+ ảnh text tiếng Việt", "60+ font chữ thông dụng"])

    # ── 10. KIE ────────────────────────────────────────────
    s = slide(prs, "KEY INFORMATION EXTRACTION — Từ rule-based đến Multi-modal")
    box(s, Inches(0.3), Inches(1.2), Inches(4), Inches(2.5), LIGHT_BG, GRAY_TEXT,
        "Thế hệ 1 — Rule-based", ["Regex, heuristic", "Nhanh, không cần training", "Vỡ ngay khi gặp format mới"])
    box(s, Inches(0.3), Inches(3.9), Inches(4), Inches(2.5), RGBColor(0xFF,0xF3,0xE0), WARNING,
        "Thế hệ 2 — BERT NER", ["Token classification", "Bỏ qua thông tin vị trí không gian"])
    box(s, Inches(4.6), Inches(1.2), Inches(8.4), Inches(5.2), ACCENT_LIGHT, ACCENT,
        "★ Thế hệ 3 — LayoutLMv3 (Microsoft, 2022)", ["Kết hợp 3 nguồn: text OCR + bbox + patch ảnh",
        "Unified self-attention qua cả 3 modality", "Pre-training: MLM + MIM + Text-Image Alignment", "Đột phá cho Document AI"])

    # ── 11. Trực giác ──────────────────────────────────────

    s = slide(prs, "CON NGƯỜI ĐỌC HÓA ĐƠN NHƯ THẾ NÀO?")
    bullets(s, ["## 1. Vị trí không gian mang ngữ nghĩa:", '- "Tổng cộng" → góc dưới phải, font lớn, in đậm', "- Tên cửa hàng → trên cùng, canh giữa",
        "", "## 2. Căn lề là thông tin:", "- Giá tiền căn phải, tên mặt hàng căn trái",
        "", "## 3. Gần nhau → liên quan:", '- "2 x 29,000" và "58,000" cùng dòng → cùng mặt hàng',
        "", "## 4. Pattern lặp lại trong vùng bảng:", "- 5 dòng có cấu trúc [tên | số | số] → bảng mặt hàng"])

    # ── 12. Mô hình mô phỏng ───────────────────────────────
    s = slide(prs, "MÔ HÌNH MÔ PHỎNG TRỰC GIÁC")
    table = s.shapes.add_table(4, 2, Inches(0.5), Inches(1.5), Inches(12), Inches(4.5)).table
    for i, h in enumerate(["Trực giác", "Mô hình mô phỏng"]):
        cell = table.cell(0, i); cell.text = h; cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs: p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate([["Vùng chữ có cấu trúc pixel đặc biệt", "PaddleOCR Detection: học probability map"],
        ["Đọc ký tự cần ngữ cảnh xung quanh", 'VietOCR: attention "nhìn lại" toàn bộ ảnh'],
        ["Vị trí + nội dung + hình ảnh cùng quyết định", "LayoutLMv3: fuse 3 luồng trong transformer"]]):
        for c, val in enumerate(row):
            cell = table.cell(r+1, c); cell.text = val
            if r % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_LIGHT
            for p in cell.text_frame.paragraphs: p.font.size = Pt(13); p.font.color.rgb = DARK_TEXT

    # ── 13. DB ─────────────────────────────────────────────

    s = slide(prs, "PaddleOCR — DB (Differentiable Binarization)")
    add_img(s, img("db-algorithm.png"), Inches(7), Inches(1.3), Inches(5.8), Inches(3))
    bullets(s, ["## Phép toán chủ đạo: Segmentation + Adaptive Thresholding",
        "", "## Mô hình học 2 map song song từ feature map F:",
        "- P = σ(Conv(F))  ← Probability map", "- T = σ(Conv'(F))  ← Threshold map",
        "- B̂ = σ((P-T)/k)  ← Binary map (differentiable)",
        "", "## Bí quyết: σ((P-T)/k) với k=50 xấp xỉ hàm step nhưng vẫn có gradient",
        "", "## Loss: L = L_bce(P, Gp) + α·L_bce(B̂, Gb) + β·L1(T, Gt)",
        "## Tham số: ResNet-18/50 + FPN, khoảng 10–27M"], width=Inches(6.5))

    # ── 14. VietOCR ────────────────────────────────────────
    s = slide(prs, "VietOCR — Seq2Seq với Attention")
    add_img(s, img("vietocr-attention.png"), Inches(7), Inches(1.3), Inches(5.8), Inches(3))
    bullets(s, ["## Phép toán chủ đạo: Cross-Attention trong Transformer decoder",
        "", "## Pipeline:", "- F = CNN_Encoder(ảnh) → shape: (H/4, W/4, C)",
        "- F_flat = reshape(F) + positional_embed", "- h_t = TransformerDecoder(h_{t-1}, F_flat)",
        "- p(y_t) = Softmax(Linear(h_t))",
        "", "## Đặc biệt tiếng Việt:",
        '- Khi sinh "ề", decoder attend mạnh vào vùng có dấu phụ kép',
        "- Vocab: 130–150 ký tự Unicode tiếng Việt"], width=Inches(6.5))

    # ── 15. LayoutLMv3 ─────────────────────────────────────
    s = slide(prs, "LayoutLMv3 — Multi-modal Transformer")
    add_img(s, img("layoutlmv3.png"), Inches(7), Inches(1.3), Inches(5.8), Inches(3))
    bullets(s, ["## Phép toán chủ đạo: Unified Self-Attention qua cả 3 modality",
        "", "## Input token i:", "- e_i = E_text(token_i) + E_1D(pos_i) + E_2D(x1,y1,x2,y2,w,h)",
        "", "## Image patch j:", "- v_j = LinearProj(patch_j) + E_img_pos(j)",
        "", "## Pre-training objectives:", "- Masked Language Modeling (MLM) trên text",
        "- Masked Image Modeling (MIM) trên ảnh", "- Text-Image Alignment"], width=Inches(6.5))

    # ── 16. Loss ───────────────────────────────────────────

    s = slide(prs, "HÀM MẤT MÁT & LÝ DO LỰA CHỌN")
    table = s.shapes.add_table(4, 3, Inches(0.5), Inches(1.3), Inches(12), Inches(3)).table
    for i, h in enumerate(["Tầng", "Loss Function", "Tại sao?"]):
        cell = table.cell(0, i); cell.text = h; cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs: p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate([["Detection", "BCE + L1", "Gradient ổn định, robust với outlier"],
        ["Recognition", "Cross-Entropy chuỗi", "Smooth, vi phân được"],
        ["KIE", "Cross-Entropy/token", "Tính được tại từng token"]]):
        for c, val in enumerate(row):
            cell = table.cell(r+1, c); cell.text = val
            if r % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_LIGHT
            for p in cell.text_frame.paragraphs: p.font.size = Pt(13); p.font.color.rgb = DARK_TEXT
    bullets(s, ["## Tại sao không dùng Metric làm Loss?", "- Field Exact Match: không có gradient",
        "- CER (edit distance): không liên tục", "- Entity F1: phụ thuộc ngưỡng confidence",
        "", "→ Dùng CE làm Loss, F1/CER/Exact Match chỉ để đánh giá"], top=Inches(4.5), height=Inches(2.5))

    # ── 17. Hệ quả ────────────────────────────────────────
    s = slide(prs, "HỆ QUẢ THỰC TẾ")
    box(s, Inches(0.5), Inches(1.3), Inches(12), Inches(2), RGBColor(0xFF,0xEB,0xEE), DANGER,
        "Điểm đau lớn nhất", ["Mô hình đạt CER 98% nhưng Field Exact Match chỉ 70%",
        'Chỉ cần sai 1 ký tự trong "87.000" thành "87.500" là cả trường total sai hoàn toàn'])
    bullets(s, ["## Khả năng mở rộng Loss:", "- Thêm class mới: DISCOUNT, TAX_CODE",
        "- Weighted loss: ưu tiên TOTAL, DATE", "- Consistency constraint: Σ(subtotals) ≈ total"], top=Inches(3.8), height=Inches(3))

    # ── 18. Thách thức ─────────────────────────────────────

    s = slide(prs, "TRỞ NGẠI ĐÃ GIẢI QUYẾT & CÒN TỒN Đọng")
    box(s, Inches(0.3), Inches(1.2), Inches(6.2), Inches(5.2), RGBColor(0xE8,0xF5,0xE9), SUCCESS,
        "✓ Đã giải quyết", ["Văn bản đa hướng, cong → DB, ABCNet", "Font chữ đa dạng tiếng Việt → VietOCR pretrained",
        "Hiểu layout cơ bản → 2D positional embedding"])
    box(s, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.2), RGBColor(0xFF,0xF3,0xE0), WARNING,
        "✗ Còn tồn đọng", ["Ảnh chất lượng thấp — CER tụt 85–90%", "Hóa đơn in nhiệt bị phai",
        "Zero-shot generalization — cần fine-tune", "Dữ liệu tiếng Việt khan hiếm", "Relation Extraction phức tạp"])

    # ── 19. Tương lai ──────────────────────────────────────
    s = slide(prs, "HƯỚNG ĐI TƯƠNG LAI")
    bullets(s, ["## 1. End-to-end Document Understanding:", "- Donut, Nougat — bỏ qua OCR, ảnh → JSON trực tiếp",
        "", "## 2. Large Vision-Language Models (LVLMs):", "- GPT-4V, Gemini — zero-shot, không cần fine-tune",
        "", "## 3. Self-supervised pretraining trên tài liệu Việt:", "- Corpus hóa đơn/biểu mẫu tiếng Việt quy mô lớn",
        "", "## 4. Synthetic data generation:", "- Tạo hóa đơn giả + domain randomization"])

    # ── 20. Kết luận ───────────────────────────────────────

    s = slide(prs, "KẾT LUẬN")
    bullets(s, ["- Bài toán trích xuất hóa đơn tiếng Việt đòi hỏi giải quyết đồng thời 3 thách thức:",
        "**phát hiện vùng chữ**, **nhận dạng ký tự**, **hiểu cấu trúc ngữ nghĩa**",
        "", "- Pipeline **PaddleOCR → VietOCR → LayoutLMv3** là lựa chọn hợp lý hiện tại",
        "", "- Nhược điểm cố hữu: **lỗi tích lũy qua từng tầng**",
        "", "- Ranh giới giữa OCR+NLP pipeline và end-to-end LVLM đang dần mờ đi"])

    # ── 21. Cảm ơn ─────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, PRIMARY)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-2), Inches(7), Inches(7))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.fill.fore_color.brightness = 0.3; c.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = "CẢM ƠN THẦY VÀ CÁC BẠN!"; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "Xin cảm ơn!"; p2.font.size = Pt(24); p2.font.color.rgb = ACCENT_LIGHT; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(20)
    tb2 = s.shapes.add_textbox(Inches(1), Inches(5.5), Inches(10), Inches(1))
    tf2 = tb2.text_frame; tf2.paragraphs[0].text = "Học viện Công nghệ Bưu chính Viễn thông — Khoa CNTT2"; tf2.paragraphs[0].font.size = Pt(16); tf2.paragraphs[0].font.color.rgb = ACCENT_LIGHT; tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    out = "/home/soang/.openclaw/workspace/smartvn-report/slides/OCR-Presentation.pptx"
    prs.save(out)
    print(f"✅ Saved: {out} ({len(prs.slides)} slides)")

if __name__ == "__main__":
    main()
