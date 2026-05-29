#!/usr/bin/env python3
"""
SmartVN E-Commerce PowerPoint — v2 (compact ~20 slides, mobile-focused)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ─────────────────────────────────────────────────
PRIMARY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0x2E, 0x86, 0xDE)
ACCENT_LIGHT = RGBColor(0xD6, 0xEA, 0xF8)
SUCCESS = RGBColor(0x27, 0xAE, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
GRAY_TEXT = RGBColor(0x7F, 0x8C, 0x8D)

IMG = "/home/soang/.openclaw/workspace/smartvn-report/images"

def img(name):
    return os.path.join(IMG, name)

def add_bg(slide, color=PRIMARY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_bar(slide, color=ACCENT, top=0, height=Inches(0.06)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, Inches(13.333), height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_footer(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), Inches(13.333), Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.paragraphs[0].text = "SmartVN — Hệ thống TMĐT Microservices"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def title_text(slide, text, top=Inches(0.15), size=Pt(26)):
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = PRIMARY

def bullets(slide, items, left=Inches(0.5), top=Inches(1.1), width=Inches(6.2), height=Inches(5.5)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if item.startswith("##"):
            p.text = item[2:].strip()
            p.font.size = Pt(17)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.space_before = Pt(14)
        elif item.startswith("- "):
            p.text = "• " + item[2:]
            p.font.size = Pt(15)
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(5)
        else:
            p.text = item
            p.font.size = Pt(15)
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(5)

def add_image(slide, path, left, top, width, height=None):
    if os.path.exists(path):
        try:
            if height:
                slide.shapes.add_picture(path, left, top, width, height)
            else:
                slide.shapes.add_picture(path, left, top, width)
        except Exception as e:
            print(f"Image error: {path}: {e}")

def new_slide(prs, title, bg_color=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if bg_color:
        add_bg(slide, bg_color)
    else:
        add_bar(slide)
        add_footer(slide)
        title_text(slide, title)
    return slide

# ══════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── 1. Cover ───────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(5), Inches(5))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.fill.fore_color.brightness = 0.3; c.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "HỆ THỐNG THƯƠNG MẠI ĐIỆN TỬ\nSMARTVN"; p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.text = "Kiến trúc Microservices • Spring Boot • Mobile Android • AI Recommendation"; p2.font.size = Pt(16); p2.font.color.rgb = ACCENT_LIGHT; p2.space_before = Pt(14)
    tb2 = s.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    tf2 = tb2.text_frame; tf2.paragraphs[0].text = "Học viện Công nghệ Bưu chính Viễn thông — Khoa CNTT2\nThành phố Hồ Chí Minh, Tháng 6 năm 2026"; tf2.paragraphs[0].font.size = Pt(13); tf2.paragraphs[0].font.color.rgb = ACCENT_LIGHT

    # ── 2. Outline ─────────────────────────────────────────
    s = new_slide(prs, "NỘI DUNG")
    bullets(s, [
        "## 1. Tổng quan đề tài",
        "## 2. Kiến trúc Microservices",
        "## 3. Ứng dụng Mobile Android (MVVM)",
        "## 4. Giao diện Mobile",
        "## 5. Recommendation System",
        "## 6. Kiểm thử & Hiệu năng",
        "## 7. Kết luận",
    ])

    # ── 3. Section: Tổng quan ──────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, ACCENT)
    tb = s.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = "01  TỔNG QUAN ĐỀ TÀI"; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE

    # ── 4. Bối cảnh ────────────────────────────────────────
    s = new_slide(prs, "BỐI CẢNH & BÀI TOÁN")
    bullets(s, [
        "## Thương mại điện tử Việt Nam",
        "- Thị trường TMĐT toàn cầu > 6.000 tỷ USD (2025)",
        "- Shopee, Lazada, Tiki, Sendo cạnh tranh mạnh",
        "",
        "## Thách thức",
        "- Mở rộng linh hoạt cho lượng lớn người dùng đồng thời",
        "- Đảm bảo tính sẵn sàng — một service lỗi không sập toàn hệ thống",
        "- Bảo mật thông tin & thanh toán",
        "",
        "## Hướng tiếp cận: Microservices",
        "- Spring Boot 3.3 + Spring Cloud",
        "- Mobile Android native (Java, MVVM)",
        "- AI Recommendation (PhoBERT + ALS)",
    ])

    # ── 5. Quy trình trước/sau ─────────────────────────────
    s = new_slide(prs, "QUY TRÌNH MUA HÀNG — TRƯỚC vs SAU")
    add_image(s, img("img-014.jpg"), Inches(0.3), Inches(1.2), Inches(6.2), Inches(5.5))
    add_image(s, img("img-015.jpg"), Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.5))

    # ── 6. Section: Kiến trúc ──────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, ACCENT)
    tb = s.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = "02  KIẾN TRÚC MICROSERVICES"; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE

    # ── 7. Architecture ────────────────────────────────────
    s = new_slide(prs, "SƠ ĐỒ KIẾN TRÚC TỔNG THỂ")
    add_image(s, img("img-018.jpg"), Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5))

    # ── 8. Services table ──────────────────────────────────
    s = new_slide(prs, "CÁC MICROSERVICES")
    table = s.shapes.add_table(9, 3, Inches(0.5), Inches(1.2), Inches(12), Inches(5)).table
    headers = ["Service", "Port", "Vai trò"]
    data = [
        ["Eureka Server", "8761", "Service Discovery & Registry"],
        ["Config Server", "8888", "Quản lý cấu hình tập trung"],
        ["API Gateway", "8080", "Routing, JWT validation, CORS"],
        ["User Service", "8081", "Quản lý người dùng, JWT/OAuth2"],
        ["Product Service", "8082", "Sản phẩm + Redis cache"],
        ["Order Service", "8083", "Đơn hàng + Thanh toán VNPay"],
        ["Admin Service", "8084", "Quản trị, dashboard"],
        ["Recommend Service", "8085", "AI Recommendation (PhoBERT+ALS)"],
    ]
    for i, h in enumerate(headers):
        cell = table.cell(0, i); cell.text = h; cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs: p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r+1, c); cell.text = val
            if r % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_LIGHT
            for p in cell.text_frame.paragraphs: p.font.size = Pt(13); p.font.color.rgb = DARK_TEXT

    # ── 9. Tech stack ──────────────────────────────────────
    s = new_slide(prs, "CÔNG NGHỆ SỬ DỤNG")
    bullets(s, [
        "## Backend",
        "- Spring Boot 3.3 + Spring Cloud (Gateway, Eureka, Config)",
        "- Spring Security + JWT + OAuth2 (Google, GitHub)",
        "- Resilience4j Circuit Breaker",
        "",
        "## Data",
        "- MySQL 8.0 — RDBMS chính",
        "- Redis 7 — Cache (giảm latency 5.5x)",
        "",
        "## Frontend & Mobile",
        "- React + Vite (Customer & Admin Web)",
        "- Android Java + MVVM (Mobile)",
        "",
        "## AI & DevOps",
        "- PhoBERT + ALS Hybrid Recommendation",
        "- Docker Compose deployment",
        "- Grafana k6 performance testing",
    ])

    # ── 10. Section: Mobile ────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, ACCENT)
    tb = s.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = "03  ỨNG DỤNG MOBILE ANDROID"; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE

    # ── 11. Mobile Architecture ────────────────────────────
    s = new_slide(prs, "KIẾN TRÚC MOBILE — MVVM")
    bullets(s, [
        "## Tech Stack",
        "- Java + Android SDK",
        "- Retrofit 2 — REST API calls",
        "- Room Database — offline-first",
        "- LiveData + ViewModel — reactive UI",
        "- Material Design 3 — giao diện hiện đại",
        "",
        "## Features chính",
        "- Đăng ký / Đăng nhập (JWT + OAuth2 Google)",
        "- Duyệt sản phẩm, tìm kiếm, lọc danh mục",
        "- Giỏ hàng (offline + sync server)",
        "- Đặt hàng & thanh toán VNPay",
        "- Theo dõi trạng thái đơn hàng",
        "- Quản lý tài khoản & địa chỉ",
    ])
    add_image(s, img("img-019.jpg"), Inches(7), Inches(1.2), Inches(5.8), Inches(5.5))

    # ── 12. Mobile Use Case ────────────────────────────────
    s = new_slide(prs, "USE CASE DIAGRAM — MOBILE")
    bullets(s, [
        "## Khách vãng lai (Guest)",
        "- Xem sản phẩm, tìm kiếm, xem đánh giá",
        "",
        "## Khách hàng (Logged-in)",
        "- Tất cả quyền Guest +",
        "- Nhận gợi ý AI (PhoBERT + ALS)",
        "- Giỏ hàng → Đặt hàng → Thanh toán VNPay",
        "- Đánh giá sản phẩm",
        "- Quản lý tài khoản, địa chỉ, lịch sử đơn",
        "",
        "## Luồng chính",
        "- Browse → Add to Cart → Checkout → Payment → Track Order",
    ])
    add_image(s, img("img-019.jpg"), Inches(7), Inches(1.2), Inches(5.8), Inches(5.5))

    # ── 13. Mobile: Login UI ───────────────────────────────
    s = new_slide(prs, "MOBILE — ĐĂNG NHẬP & ĐĂNG KÝ")
    add_image(s, img("img-080.jpg"), Inches(0.5), Inches(1.3), Inches(5.5), Inches(5))
    bullets(s, [
        "## Đăng nhập",
        "- Email + Password",
        "- OAuth2 Google (one-tap)",
        "- Auto-refresh JWT token",
        "",
        "## Đăng ký",
        "- Form đăng ký + xác nhận OTP",
        "- Validate email real-time",
        "",
        "## Offline-first",
        "- Lưu token vào SharedPreferences",
        "- Auto-login khi mất mạng",
    ], left=Inches(6.5), top=Inches(1.3), width=Inches(6.3))

    # ── 14. Mobile: Product UI ─────────────────────────────
    s = new_slide(prs, "MOBILE — DANH SÁCH & CHI TIẾT SẢN PHẨM")
    add_image(s, img("img-085.jpg"), Inches(0.5), Inches(1.3), Inches(5.5), Inches(5))
    bullets(s, [
        "## Danh sách sản phẩm",
        "- RecyclerView + GridLayoutManager",
        "- Filter theo danh mục (Smartphones, Laptop...)",
        "- Price range slider",
        "- Pull-to-refresh",
        "",
        "## Chi tiết sản phẩm",
        "- ViewPager2 cho gallery ảnh",
        "- Thông số kỹ thuật (RAM, pin, màn hình...)",
        "- Nút Add to Cart + Buy Now",
        "- Gợi ý sản phẩm tương tự (AI)",
    ], left=Inches(6.5), top=Inches(1.3), width=Inches(6.3))

    # ── 15. Mobile: Cart & Checkout ────────────────────────
    s = new_slide(prs, "MOBILE — GIỎ HÀNG & CHECKOUT")
    bullets(s, [
        "## Giỏ hàng",
        "- Room Database lưu offline",
        "- Tăng/giảm số lượng, xoá sản phẩm",
        "- Sync với server khi có mạng",
        "",
        "## Checkout",
        "- Chọn địa chỉ giao hàng",
        "- Chọn phương thức thanh toán",
        "- Xác nhận đơn hàng",
        "",
        "## Thanh toán VNPay",
        "- WebView redirect đến cổng VNPay",
        "- Callback xử lý kết quả",
        "- Cập nhật trạng thái đơn hàng",
    ])
    add_image(s, img("img-095.jpg"), Inches(7), Inches(1.3), Inches(5.8), Inches(5))

    # ── 16. Mobile: Profile & Orders ───────────────────────
    s = new_slide(prs, "MOBILE — TÀI KHOẢN & ĐƠN HÀNG")
    add_image(s, img("img-090.jpg"), Inches(0.5), Inches(1.3), Inches(5.5), Inches(5))
    bullets(s, [
        "## Quản lý tài khoản",
        "- Profile Info: tên, email, SĐT",
        "- Quản lý địa chỉ (thêm/sửa/xoá)",
        "- Đổi mật khẩu",
        "",
        "## Đơn hàng",
        "- Danh sách đơn + trạng thái",
        "- Chi tiết đơn (sản phẩm, tổng tiền)",
        "- Theo dõi giao hàng",
        "- Huỷ đơn (nếu chưa ship)",
        "",
        "## Đánh giá",
        "- Rate sản phẩm (1-5 sao)",
        "- Viết review + ảnh",
    ], left=Inches(6.5), top=Inches(1.3), width=Inches(6.3))

    # ── 17. Section: Recommendation ────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, ACCENT)
    tb = s.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = "04  AI RECOMMENDATION SYSTEM"; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE

    # ── 18. Recommendation ─────────────────────────────────
    s = new_slide(prs, "HYBRID RECOMMENDATION — PhoBERT + ALS")
    bullets(s, [
        "## Content-Based: PhoBERT",
        "- Pre-trained Vietnamese language model (VinAI)",
        "- Phân tích mô tả sản phẩm → text embeddings",
        "- Tính similarity giữa sản phẩm",
        "",
        "## Collaborative: ALS",
        "- Alternating Least Squares",
        "- Học từ lịch sử mua hàng user-item matrix",
        "- Dự đoán user sẽ thích sản phẩm nào",
        "",
        "## Hybrid Merge",
        "- Kết hợp điểm từ cả 2 model",
        "- Cold start: user mới → Content-Based",
        "- Warm user → Collaborative ưu tiên hơn",
    ])
    add_image(s, img("img-040.jpg"), Inches(7), Inches(1.2), Inches(5.8), Inches(5.5))

    # ── 19. Section: Testing ───────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, ACCENT)
    tb = s.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = "05  KIỂM THỬ & HIỆU NĂNG"; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE

    # ── 20. Benchmark ──────────────────────────────────────
    s = new_slide(prs, "KẾT QUẢ BENCHMARK — REDIS CACHE")
    table = s.shapes.add_table(5, 4, Inches(0.5), Inches(1.5), Inches(12), Inches(4)).table
    headers = ["Metric", "Không Cache", "Redis Cache", "Cải thiện"]
    data = [
        ["Avg Response Time", "47.81ms", "8.66ms", "5.5x nhanh hơn"],
        ["Throughput", "196.8 req/s", "248.5 req/s", "+26%"],
        ["Error Rate", "> 0%", "0%", "Hoàn hảo"],
        ["P95 Latency", "~120ms", "~25ms", "4.8x nhanh hơn"],
    ]
    for i, h in enumerate(headers):
        cell = table.cell(0, i); cell.text = h; cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs: p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r+1, c); cell.text = val
            if r % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_LIGHT
            for p in cell.text_frame.paragraphs: p.font.size = Pt(14); p.font.color.rgb = DARK_TEXT
    bullets(s, [
        "## Tool: Grafana k6",
        "- Load testing với 100 virtual users",
        "- So sánh: No Cache vs Redis Cache",
        "- Product Service endpoint chính",
    ], top=Inches(5.5), height=Inches(1.5))

    # ── 21. Kết luận ───────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, ACCENT)
    tb = s.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = "06  KẾT LUẬN"; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = WHITE

    # ── 22. Kết quả ────────────────────────────────────────
    s = new_slide(prs, "KẾT QUẢ & HƯỚNG PHÁT TRIỂN")
    bullets(s, [
        "## Kết quả đạt được",
        "- Hệ thống TMĐT SmartVN hoàn chỉnh (7 microservices)",
        "- Ứng dụng Mobile Android native (MVVM) đầy đủ tính năng",
        "- AI Recommendation (PhoBERT + ALS Hybrid)",
        "- Hiệu năng: Redis cache giảm latency 5.5x",
        "- Bảo mật: JWT + OAuth2 + Circuit Breaker",
        "- Docker Compose deployment",
        "",
        "## Hướng phát triển",
        "- CI/CD pipeline (GitHub Actions)",
        "- Monitoring (Prometheus + Grafana)",
        "- Push notification (Firebase FCM)",
        "- Elasticsearch full-text search",
        "- iOS app (Flutter cross-platform)",
    ])

    # ── 23. Cảm ơn ─────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-2), Inches(7), Inches(7))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.fill.fore_color.brightness = 0.3; c.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = "CẢM ƠN THẦY VÀ CÁC BẠN!"; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "Xin cảm ơn!"; p2.font.size = Pt(24); p2.font.color.rgb = ACCENT_LIGHT; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(20)
    tb2 = s.shapes.add_textbox(Inches(1), Inches(5.5), Inches(10), Inches(1))
    tf2 = tb2.text_frame; tf2.paragraphs[0].text = "Học viện Công nghệ Bưu chính Viễn thông — Khoa CNTT2"; tf2.paragraphs[0].font.size = Pt(16); tf2.paragraphs[0].font.color.rgb = ACCENT_LIGHT; tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ── Save ───────────────────────────────────────────────
    out = "/home/soang/.openclaw/workspace/smartvn-report/SmartVN-Slides.pptx"
    prs.save(out)
    print(f"✅ Saved: {out} ({len(prs.slides)} slides)")

if __name__ == "__main__":
    main()
