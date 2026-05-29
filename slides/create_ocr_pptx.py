#!/usr/bin/env python3
"""
OCR Presentation → PPTX with animations & transitions
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import lxml.etree as etree

# ── Colors ─────────────────────────────────────────────────
PRIMARY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0x2E, 0x86, 0xDE)
ACCENT_LIGHT = RGBColor(0xD6, 0xEA, 0xF8)
ACCENT_DARK = RGBColor(0x1A, 0x52, 0x76)
SUCCESS = RGBColor(0x27, 0xAE, 0x60)
WARNING = RGBColor(0xF3, 0x9C, 0x12)
DANGER = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
GRAY_TEXT = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)

def add_transition(slide, transition_type="fade", duration=1500):
    """Add slide transition with animation."""
    transition = slide._element.find(qn('p:transition'))
    if transition is None:
        transition = etree.SubElement(slide._element, qn('p:transition'))
    transition.set('spd', 'med')
    transition.set('advTm', str(duration))
    
    # Add transition type
    if transition_type == "fade":
        etree.SubElement(transition, qn('p:fade'))
    elif transition_type == "push":
        etree.SubElement(transition, qn('p:push'))
    elif transition_type == "wipe":
        etree.SubElement(transition, qn('p:wipe'))
    elif transition_type == "split":
        etree.SubElement(transition, qn('p:split'))
    elif transition_type == "zoom":
        etree.SubElement(transition, qn('p:zoom'))

def add_appear_animation(shape, delay=0):
    """Add appear animation to a shape."""
    sp = shape._element
    sp_id = sp.get('id', '1')
    
    # Build animation XML
    timing = etree.SubElement(sp, qn('p:timing'))
    tnLst = etree.SubElement(timing, qn('p:tnLst'))
    
    par = etree.SubElement(tnLst, qn('p:par'))
    cTn = etree.SubElement(par, qn('p:cTn'))
    cTn.set('id', '1')
    cTn.set('dur', 'indefinite')
    cTn.set('restart', 'never')
    cTn.set('nodeType', 'tmRoot')
    
    childTnLst = etree.SubElement(cTn, qn('p:childTnLst'))
    seq = etree.SubElement(childTnLst, qn('p:seq'))
    seq.set('concurrent', '1')
    seq.set('nextAc', 'seek')
    
    seqCTn = etree.SubElement(seq, qn('p:cTn'))
    seqCTn.set('id', '2')
    seqCTn.set('dur', 'indefinite')
    seqCTn.set('nodeType', 'mainSeq')
    
    seqChildTnLst = etree.SubElement(seqCTn, qn('p:childTnLst'))

def add_bg(slide, color=PRIMARY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_bar(slide, color=ACCENT, top=0, height=Inches(0.06)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, Inches(13.333), height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_footer(slide, text="Trích xuất Hóa đơn Tiếng Việt — Pipeline PaddleOCR • VietOCR • LayoutLMv3"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), Inches(13.333), Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def title_text(slide, text, top=Inches(0.15), size=Pt(26)):
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = PRIMARY

def bullets(slide, items, left=Inches(0.5), top=Inches(1.1), width=Inches(12), height=Inches(5.5)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if item.startswith("##"):
            p.text = item[2:].strip()
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.space_before = Pt(14)
        elif item.startswith("!"):
            p.text = "⚠ " + item[1:].strip()
            p.font.size = Pt(15)
            p.font.color.rgb = DANGER
            p.space_before = Pt(6)
        elif item.startswith("- "):
            p.text = "• " + item[2:]
            p.font.size = Pt(15)
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(5)
        else:
            p.text = item
            p.font.size = Pt(15)
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(5)

def add_colored_box(slide, left, top, width, height, fill_color, border_color, title, content_lines):
    """Add a colored box with title and content."""
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.4))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = border_color
    title_bar.line.fill.background()
    tf = title_bar.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Content box
    content_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top + Inches(0.4), width, height - Inches(0.4))
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = fill_color
    content_box.line.color.rgb = border_color
    content_box.line.width = Pt(1)
    tf2 = content_box.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Pt(10)
    tf2.margin_right = Pt(10)
    tf2.margin_top = Pt(5)
    
    for i, line in enumerate(content_lines):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = "• " + line
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(3)

def new_slide(prs, title, bg_color=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if bg_color:
        add_bg(slide, bg_color)
    else:
        add_bar(slide)
        add_footer(slide)
        title_text(slide, title)
    return slide

def add_section(prs, num, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, ACCENT)
    # Number
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(2), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    # Title
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(10), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    add_transition(slide, "fade", 1000)
    return slide

# ══════════════════════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ── 1. Title Slide ─────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    # Decorative circle
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(5), Inches(5))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.fill.fore_color.brightness = 0.3; c.line.fill.background()
    # Title
    tb = s.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "TRÍCH XUẤT THÔNG TIN\nHÓA ĐƠN TIẾNG VIỆT TỪ ẢNH"
    p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = "Pipeline PaddleOCR — VietOCR — LayoutLMv3"
    p2.font.size = Pt(18); p2.font.color.rgb = ACCENT_LIGHT; p2.space_before = Pt(14)
    p3 = tf.add_paragraph()
    p3.text = "Đồ án Môn học"
    p3.font.size = Pt(16); p3.font.color.rgb = GRAY_TEXT; p3.space_before = Pt(20)
    # Institute
    tb2 = s.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    tf2 = tb2.text_frame
    tf2.paragraphs[0].text = "Học viện Công nghệ Bưu chính Viễn thông\nKhoa Công nghệ Thông tin 2 — Tháng 6, 2026"
    tf2.paragraphs[0].font.size = Pt(14); tf2.paragraphs[0].font.color.rgb = ACCENT_LIGHT
    add_transition(s, "fade", 2000)
    
    # ── 2. Outline ─────────────────────────────────────────
    s = new_slide(prs, "NỘI DUNG TRÌNH BÀY")
    bullets(s, [
        "## 1. Mở đầu — Bối cảnh bài toán",
        "## 2. Tổng quan — Mô hình hóa & Pipeline",
        "## 3. Phân loại & Tiến hóa — 3 thế hệ thuật toán",
        "## 4. Trực giác — Con người đọc hóa đơn ra sao?",
        "## 5. Toán học hóa — Công thức cốt lõi",
        "## 6. Khoảng cách — Lý thuyết vs Thực tế",
        "## 7. Thách thức — Rào cản & Hướng phát triển",
        "## 8. Kết luận",
    ])
    add_transition(s, "fade", 800)
    
    # ── Section 1: MỞ ĐẦU ─────────────────────────────────
    add_section(prs, "01", "MỞ ĐẦU")
    
    # Slide 3: Bối cảnh
    s = new_slide(prs, "BỐI CẢNH BÀI TOÁN")
    bullets(s, [
        "- Mỗi ngày, hàng triệu giao dịch thương mại tại Việt Nam tạo ra hàng triệu tờ hóa đơn giấy",
        "- Kế toán viên gõ tay từng con số → sai sót nhập liệu thường xuyên",
        "",
        "## Thách thức đọc hóa đơn thực tế:",
        "- Font chữ tùy ý, góc chụp nghiêng",
        "- Ánh sáng không đều",
        "- Tiếng Việt với đầy đủ dấu thanh",
        "",
        "- Cần biến **ảnh hóa đơn** thành **dữ liệu có cấu trúc**",
    ])
    add_transition(s, "fade", 800)
    
    # Slide 4: Tình huống thực tế
    s = new_slide(prs, "CÁC TÌNH HUỐNG THỰC TẾ")
    add_colored_box(s, Inches(0.5), Inches(1.3), Inches(12), Inches(1.3),
        ACCENT_LIGHT, ACCENT, "Kịch bản A — Kế toán doanh nghiệp nhỏ",
        ["300 hóa đơn công tác phí, mất 2–3 ngày nhập thủ công"])
    add_colored_box(s, Inches(0.5), Inches(2.8), Inches(12), Inches(1.3),
        RGBColor(0xE8, 0xF5, 0xE9), SUCCESS, "Kịch bản B — Ứng dụng chi tiêu cá nhân",
        ["Chụp ảnh hóa đơn → tự động điền ngày, siêu thị, tổng tiền"])
    add_colored_box(s, Inches(0.5), Inches(4.3), Inches(12), Inches(1.3),
        RGBColor(0xFF, 0xF3, 0xE0), WARNING, "Kịch bản C — Kiểm toán & Tuân thủ thuế",
        ["Xác minh hàng nghìn hóa đơn VAT tự động"])
    # Arrow text
    txBox = s.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "→ Nhu cầu chung: Document Information Extraction (DIE)"
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = ACCENT; p.alignment = PP_ALIGN.CENTER
    add_transition(s, "fade", 800)
    
    # ── Section 2: TỔNG QUAN ───────────────────────────────
    add_section(prs, "02", "TỔNG QUAN BÀI TOÁN")
    
    # Slide 5: Ba tầng xử lý
    s = new_slide(prs, "MÔ HÌNH HÓA — BA TẦNG XỬ LÝ")
    # Tier 3
    add_colored_box(s, Inches(1.5), Inches(1.3), Inches(10), Inches(1.5),
        ACCENT_LIGHT, ACCENT, "Tầng 3: Trích xuất thực thể có cấu trúc (KIE)",
        ["LayoutLMv3 — Multi-modal Transformer"])
    # Arrow
    txBox = s.shapes.add_textbox(Inches(6), Inches(2.8), Inches(1), Inches(0.5))
    tf = txBox.text_frame; tf.paragraphs[0].text = "↓"; tf.paragraphs[0].font.size = Pt(24); tf.paragraphs[0].font.color.rgb = ACCENT; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Tier 2
    add_colored_box(s, Inches(1.5), Inches(3.2), Inches(10), Inches(1.5),
        RGBColor(0xE3, 0xF2, 0xFD), ACCENT_DARK, "Tầng 2: Nhận dạng ký tự (OCR Recognition)",
        ["VietOCR — Seq2Seq + Attention"])
    # Arrow
    txBox = s.shapes.add_textbox(Inches(6), Inches(4.7), Inches(1), Inches(0.5))
    tf = txBox.text_frame; tf.paragraphs[0].text = "↓"; tf.paragraphs[0].font.size = Pt(24); tf.paragraphs[0].font.color.rgb = ACCENT; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Tier 1
    add_colored_box(s, Inches(1.5), Inches(5.1), Inches(10), Inches(1.5),
        LIGHT_BG, PRIMARY, "Tầng 1: Phát hiện vùng chữ (Text Detection)",
        ["PaddleOCR — DB (Differentiable Binarization)"])
    add_transition(s, "fade", 800)
    
    # Slide 6: Input/Output
    s = new_slide(prs, "ĐẦU VÀO & ĐẦU RA")
    add_colored_box(s, Inches(0.5), Inches(1.3), Inches(12), Inches(1.2),
        ACCENT_LIGHT, ACCENT, "Đầu vào",
        ["Ảnh JPEG/PNG hóa đơn chụp bằng điện thoại (nghiêng, nhòe, ánh sáng không đều)"])
    add_colored_box(s, Inches(0.5), Inches(3), Inches(12), Inches(3.5),
        RGBColor(0xE8, 0xF5, 0xE9), SUCCESS, "Đầu ra — JSON có ngữ nghĩa",
        ['{ "store_name": "Circle K", "date": "2024-05-27", "total": 87000',
         '  "items": [ { "name": "Cafe sua da", "qty": 2, "price": 29000 }, ... ] }'])
    add_transition(s, "fade", 800)
    
    # Slide 7: Pipeline
    s = new_slide(prs, "PIPELINE XỬ LÝ")
    # Boxes
    box_data = [
        ("Ảnh hóa đơn", LIGHT_BG, PRIMARY),
        ("PaddleOCR\nDetection", ACCENT_LIGHT, ACCENT),
        ("Bounding\nBoxes", LIGHT_BG, PRIMARY),
        ("VietOCR\nRecognition", ACCENT_LIGHT, ACCENT),
        ("Text + BBox", LIGHT_BG, PRIMARY),
        ("LayoutLMv3\nKIE", ACCENT_LIGHT, ACCENT),
        ("JSON\nCấu trúc", RGBColor(0xE8, 0xF5, 0xE9), SUCCESS),
    ]
    for i, (text, fill, border) in enumerate(box_data):
        left = Inches(0.3 + i * 1.8)
        top = Inches(2.5) if i < 5 else Inches(5)
        left_adj = left if i < 5 else Inches(0.3 + (6-i) * 1.8)
        if i == 5: left_adj = Inches(4)
        if i == 6: left_adj = Inches(0.3)
        
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_adj, top, Inches(1.6), Inches(1))
        box.fill.solid(); box.fill.fore_color.rgb = fill; box.line.color.rgb = border
        tf = box.text_frame; tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = DARK_TEXT; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
    
    # Arrows
    for i in range(4):
        txBox = s.shapes.add_textbox(Inches(1.9 + i * 1.8), Inches(2.7), Inches(0.3), Inches(0.5))
        tf = txBox.text_frame; tf.paragraphs[0].text = "→"; tf.paragraphs[0].font.size = Pt(18); tf.paragraphs[0].font.color.rgb = ACCENT
    # Down arrow
    txBox = s.shapes.add_textbox(Inches(5.5), Inches(3.7), Inches(1), Inches(0.5))
    tf = txBox.text_frame; tf.paragraphs[0].text = "↓"; tf.paragraphs[0].font.size = Pt(24); tf.paragraphs[0].font.color.rgb = ACCENT; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Left arrow
    txBox = s.shapes.add_textbox(Inches(3.5), Inches(5.2), Inches(0.5), Inches(0.5))
    tf = txBox.text_frame; tf.paragraphs[0].text = "←"; tf.paragraphs[0].font.size = Pt(18); tf.paragraphs[0].font.color.rgb = ACCENT
    add_transition(s, "wipe", 800)
    
    # ── Section 3: PHÂN LOẠI ───────────────────────────────
    add_section(prs, "03", "PHÂN LOẠI & TIẾN HÓA")
    
    # Slide 8: Text Detection
    s = new_slide(prs, "TEXT DETECTION — Từ sliding window đến anchor-free")
    add_colored_box(s, Inches(0.3), Inches(1.2), Inches(4), Inches(2.5),
        RGBColor(0xF5, 0xF5, 0xF5), GRAY_TEXT, "Thế hệ 1 (trước 2015)",
        ["MSER, SWT", "Phân tích pixel theo màu sắc", "Chỉ hoạt động với ảnh sạch"])
    add_colored_box(s, Inches(0.3), Inches(3.9), Inches(4), Inches(2.5),
        RGBColor(0xFF, 0xF3, 0xE0), WARNING, "Thế hệ 2 (2015–2017)",
        ["CTPN, TextBoxes++, SegLink", "CNN-based với Anchor Boxes", "Khó phát hiện văn bản cong"])
    add_colored_box(s, Inches(4.6), Inches(1.2), Inches(8.4), Inches(5.2),
        ACCENT_LIGHT, ACCENT, "★ Thế hệ 3 (2018–nay) — DB",
        ["DB (Differentiable Binarization)", "Xương sống của PaddleOCR",
         "Dự đoán probability map per pixel", "Ngưỡng có thể học được (differentiable threshold)",
         "Tối ưu cho CPU/mobile"])
    add_transition(s, "fade", 800)
    
    # Slide 9: Text Recognition
    s = new_slide(prs, "TEXT RECOGNITION — Từ template matching đến Attention")
    add_colored_box(s, Inches(0.3), Inches(1.2), Inches(4), Inches(2.5),
        RGBColor(0xF5, 0xF5, 0xF5), GRAY_TEXT, "Thế hệ 1 (trước 2014)",
        ["Tesseract v3, template matching", "Tách từng ký tự", "Yêu cầu font rõ ràng"])
    add_colored_box(s, Inches(0.3), Inches(3.9), Inches(4), Inches(2.5),
        RGBColor(0xFF, 0xF3, 0xE0), WARNING, "Thế hệ 2 (2015–2018)",
        ["CRNN + CTC", "CNN + BiLSTM + CTC loss", "Nền tảng OCR hiện đại"])
    add_colored_box(s, Inches(4.6), Inches(1.2), Inches(8.4), Inches(5.2),
        ACCENT_LIGHT, ACCENT, "★ Thế hệ 3 (2018–nay) — VietOCR",
        ["Seq2Seq + Attention", "VGG/ResNet encoder + Transformer decoder",
         "Xử lý dấu thanh tiếng Việt tốt hơn", "Pretrain trên 5M+ ảnh text tiếng Việt",
         "60+ font chữ thông dụng"])
    add_transition(s, "fade", 800)
    
    # Slide 10: KIE
    s = new_slide(prs, "KEY INFORMATION EXTRACTION — Từ rule-based đến Multi-modal")
    add_colored_box(s, Inches(0.3), Inches(1.2), Inches(4), Inches(2.5),
        RGBColor(0xF5, 0xF5, 0xF5), GRAY_TEXT, "Thế hệ 1 — Rule-based",
        ["Regex, heuristic", "Nhanh, không cần training", "Vỡ ngay khi gặp format mới"])
    add_colored_box(s, Inches(0.3), Inches(3.9), Inches(4), Inches(2.5),
        RGBColor(0xFF, 0xF3, 0xE0), WARNING, "Thế hệ 2 — BERT NER",
        ["Token classification", "Bỏ qua thông tin vị trí không gian"])
    add_colored_box(s, Inches(4.6), Inches(1.2), Inches(8.4), Inches(5.2),
        ACCENT_LIGHT, ACCENT, "★ Thế hệ 3 — LayoutLMv3 (Microsoft, 2022)",
        ["Kết hợp 3 nguồn: text OCR + bbox + patch ảnh",
         "Unified self-attention qua cả 3 modality",
         "Pre-training: MLM + MIM + Text-Image Alignment",
         "Đột phá cho Document AI"])
    add_transition(s, "fade", 800)
    
    # ── Section 4: TRỰC GIÁC ───────────────────────────────
    add_section(prs, "04", "TRỰC GIÁC & TƯ DUY")
    
    # Slide 11
    s = new_slide(prs, "CON NGƯỜI ĐỌC HÓA ĐƠN NHƯ THẾ NÀO?")
    bullets(s, [
        "## 1. Vị trí không gian mang ngữ nghĩa:",
        '- "Tổng cộng" → góc dưới phải, font lớn, in đậm',
        "- Tên cửa hàng → trên cùng, canh giữa",
        "",
        "## 2. Căn lề là thông tin:",
        "- Giá tiền căn phải, tên mặt hàng căn trái",
        "",
        "## 3. Gần nhau → liên quan:",
        '- "2 x 29,000" và "58,000" cùng dòng → cùng mặt hàng',
        "",
        "## 4. Pattern lặp lại trong vùng bảng:",
        "- 5 dòng có cấu trúc [tên | số | số] → bảng mặt hàng",
    ])
    add_transition(s, "fade", 800)
    
    # Slide 12
    s = new_slide(prs, "MÔ HÌNH MÔ PHỎNG TRỰC GIÁC CON NGƯỜI")
    # Table
    table = s.shapes.add_table(4, 2, Inches(0.5), Inches(1.5), Inches(12), Inches(4.5)).table
    headers = ["Trực giác", "Mô hình mô phỏng"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i); cell.text = h; cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs: p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    rows = [
        ["Vùng chữ có cấu trúc pixel đặc biệt", "PaddleOCR Detection: học probability map"],
        ["Đọc ký tự cần ngữ cảnh xung quanh", 'VietOCR: attention "nhìn lại" toàn bộ ảnh'],
        ["Vị trí + nội dung + hình ảnh cùng quyết định", "LayoutLMv3: fuse 3 luồng trong transformer"],
    ]
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r+1, c); cell.text = val
            if r % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_LIGHT
            for p in cell.text_frame.paragraphs: p.font.size = Pt(13); p.font.color.rgb = DARK_TEXT
    # Warning box
    add_colored_box(s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.8),
        RGBColor(0xFF, 0xF3, 0xE0), WARNING, "Giả định ngầm",
        ["Mỗi mô hình đều có giả định riêng → khi giả định vi phạm, mô hình suy giảm chất lượng"])
    add_transition(s, "fade", 800)
    
    # ── Section 5: TOÁN HỌC ────────────────────────────────
    add_section(prs, "05", "TOÁN HỌC HÓA")
    
    # Slide 13
    s = new_slide(prs, "PaddleOCR — DB (Differentiable Binarization)")
    add_colored_box(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.8),
        ACCENT_LIGHT, ACCENT, "Phép toán chủ đạo: Segmentation + Adaptive Thresholding", [])
    bullets(s, [
        "## Mô hình học 2 map song song từ feature map F:",
        "- P = σ(Conv(F))  ← Probability map",
        "- T = σ(Conv'(F))  ← Threshold map",
        "- B̂ = σ((P - T) / k)  ← Binary map (differentiable)",
        "",
        "## Bí quyết:",
        "- σ((P-T)/k) với k=50 xấp xỉ hàm step nhưng vẫn có gradient để backprop",
        "",
        "## Loss:",
        "- L = L_bce(P, Gp) + α·L_bce(B̂, Gb) + β·L1(T, Gt)",
        "",
        "## Tham số: ResNet-18/50 + FPN, khoảng 10–27M tham số",
    ], top=Inches(2.2))
    add_transition(s, "fade", 800)
    
    # Slide 14
    s = new_slide(prs, "VietOCR — Seq2Seq với Attention")
    add_colored_box(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.8),
        ACCENT_LIGHT, ACCENT, "Phép toán chủ đạo: Cross-Attention trong Transformer decoder", [])
    bullets(s, [
        "## Pipeline:",
        "- F = CNN_Encoder(ảnh)  → shape: (H/4, W/4, C)",
        "- F_flat = reshape(F) + positional_embed",
        "- h_t = TransformerDecoder(h_{t-1}, F_flat)",
        "- p(y_t) = Softmax(Linear(h_t))",
        "",
        "## Attention:",
        "- α_{t,s} = softmax(Q_t · K_s / √d_k) — cho phép 'nhìn' vào vùng feature tương ứng",
        "",
        "## Đặc biệt tiếng Việt:",
        '- Khi sinh "ề", decoder attend mạnh vào vùng có dấu phụ kép (mũ + huyền)',
        "- Vocab: 130–150 ký tự Unicode tiếng Việt + dấu câu + số",
    ], top=Inches(2.2))
    add_transition(s, "fade", 800)
    
    # Slide 15
    s = new_slide(prs, "LayoutLMv3 — Multi-modal Transformer")
    add_colored_box(s, Inches(0.5), Inches(1.2), Inches(12), Inches(0.8),
        ACCENT_LIGHT, ACCENT, "Phép toán chủ đạo: Unified Self-Attention qua cả 3 modality", [])
    bullets(s, [
        "## Input token i:",
        "- e_i = E_text(token_i) + E_1D(pos_i) + E_2D(x1,y1,x2,y2,w,h)",
        "",
        "## Image patch j:",
        "- v_j = LinearProj(patch_j) + E_img_pos(j)",
        "",
        "## Concatenate → Transformer → Contextualized representations",
        "",
        "## Pre-training objectives:",
        "- Masked Language Modeling (MLM) trên text",
        "- Masked Image Modeling (MIM) trên ảnh",
        "- Text-Image Alignment",
    ], top=Inches(2.2))
    add_transition(s, "fade", 800)
    
    # ── Section 6: KHOẢNG CÁCH ─────────────────────────────
    add_section(prs, "06", "KHOẢNG CÁCH GIỮA LÝ THUYẾT VÀ THỰC TẾ")
    
    # Slide 16
    s = new_slide(prs, "HÀM MẤT MÁT & LÝ DO LỰA CHỌN")
    table = s.shapes.add_table(4, 3, Inches(0.5), Inches(1.3), Inches(12), Inches(3)).table
    headers = ["Tầng", "Loss Function", "Tại sao?"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i); cell.text = h; cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs: p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    rows = [
        ["Detection", "BCE + L1", "Gradient ổn định, robust với outlier"],
        ["Recognition", "Cross-Entropy chuỗi", "Smooth, vi phân được"],
        ["KIE", "Cross-Entropy/token", "Tính được tại từng token"],
    ]
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r+1, c); cell.text = val
            if r % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT_LIGHT
            for p in cell.text_frame.paragraphs: p.font.size = Pt(13); p.font.color.rgb = DARK_TEXT
    bullets(s, [
        "## Tại sao không dùng Metric làm Loss?",
        "!Field Exact Match: không có gradient",
        "!CER (edit distance): không liên tục",
        "!Entity F1: phụ thuộc ngưỡng confidence",
        "",
        "→ Dùng CE làm Loss, F1/CER/Exact Match chỉ để đánh giá",
    ], top=Inches(4.5), height=Inches(2.5))
    add_transition(s, "fade", 800)
    
    # Slide 17
    s = new_slide(prs, "HỆ QUẢ THỰC TẾ")
    add_colored_box(s, Inches(0.5), Inches(1.3), Inches(12), Inches(2),
        RGBColor(0xFF, 0xEB, 0xEE), DANGER, "⚠ Điểm đau lớn nhất",
        ["Mô hình có thể đạt CER rất thấp (98% ký tự đúng) nhưng Field Exact Match chỉ 70%",
         "Chỉ cần sai 1 ký tự trong '87.000' thành '87.500' là cả trường total sai hoàn toàn"])
    bullets(s, [
        "## Khả năng mở rộng Loss:",
        "- Thêm class mới: DISCOUNT, TAX_CODE",
        "- Weighted loss: ưu tiên TOTAL, DATE",
        "- Consistency constraint: Σ(subtotals) ≈ total",
    ], top=Inches(3.8), height=Inches(3))
    add_transition(s, "fade", 800)
    
    # ── Section 7: THÁCH THỨC ──────────────────────────────
    add_section(prs, "07", "THÁCH THỨC & RÀO CẢN")
    
    # Slide 18
    s = new_slide(prs, "TRỞ NGẠI ĐÃ GIẢI QUYẾT & CÒN TỒN Đọng")
    add_colored_box(s, Inches(0.3), Inches(1.2), Inches(6.2), Inches(5.2),
        RGBColor(0xE8, 0xF5, 0xE9), SUCCESS, "✓ Đã giải quyết",
        ["Văn bản đa hướng, cong → DB, ABCNet (bezier curves)",
         "Font chữ đa dạng tiếng Việt → VietOCR pretrained + augmentation",
         "Hiểu layout cơ bản → 2D positional embedding (LayoutLM)"])
    add_colored_box(s, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.2),
        RGBColor(0xFF, 0xF3, 0xE0), WARNING, "✗ Còn tồn đọng",
        ["Ảnh chất lượng thấp (run, nhòe, phản chiếu) — CER tụt 85–90%",
         "Hóa đơn in nhiệt bị phai — contrast rất thấp",
         "Zero-shot generalization — cần fine-tune cho mỗi template mới",
         "Dữ liệu tiếng Việt khan hiếm, hạ tầng GPU đắt",
         "Relation Extraction phức tạp (20+ mặt hàng)"])
    add_transition(s, "fade", 800)
    
    # Slide 19
    s = new_slide(prs, "HƯỚNG ĐI TƯƠNG LAI")
    bullets(s, [
        "## 1. End-to-end Document Understanding:",
        "- Donut, Nougat — bỏ qua OCR, ảnh → JSON trực tiếp",
        "",
        "## 2. Large Vision-Language Models (LVLMs):",
        "- GPT-4V, Gemini — zero-shot, không cần fine-tune",
        "",
        "## 3. Self-supervised pretraining trên tài liệu Việt:",
        "- Corpus hóa đơn/biểu mẫu tiếng Việt quy mô lớn",
        "",
        "## 4. Synthetic data generation:",
        "- Tạo hóa đơn giả + domain randomization",
    ])
    add_transition(s, "fade", 800)
    
    # ── Section 8: KẾT LUẬN ────────────────────────────────
    add_section(prs, "08", "KẾT LUẬN")
    
    # Slide 20
    s = new_slide(prs, "KẾT LUẬN")
    bullets(s, [
        "- Bài toán trích xuất hóa đơn tiếng Việt đòi hỏi giải quyết đồng thời 3 thách thức: **phát hiện vùng chữ**, **nhận dạng ký tự**, **hiểu cấu trúc ngữ nghĩa**",
        "",
        "- Pipeline **PaddleOCR → VietOCR → LayoutLMv3** là lựa chọn hợp lý hiện tại với pretrained weights chất lượng tốt",
        "",
        "- Nhược điểm cố hữu: **lỗi tích lũy qua từng tầng**",
        "",
        "- Ranh giới giữa OCR+NLP pipeline và end-to-end LVLM đang dần mờ đi",
    ])
    add_transition(s, "fade", 800)
    
    # ── Slide 21: Cảm ơn ───────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PRIMARY)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-2), Inches(7), Inches(7))
    c.fill.solid(); c.fill.fore_color.rgb = ACCENT; c.fill.fore_color.brightness = 0.3; c.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.text = "CẢM ƠN THẦY VÀ CÁC BẠN!"; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = "Xin cảm ơn!"; p2.font.size = Pt(24); p2.font.color.rgb = ACCENT_LIGHT; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(20)
    tb2 = s.shapes.add_textbox(Inches(1), Inches(5.5), Inches(10), Inches(1))
    tf2 = tb2.text_frame; tf2.paragraphs[0].text = "Học viện Công nghệ Bưu chính Viễn thông — Khoa CNTT2\nTháng 6, 2026"; tf2.paragraphs[0].font.size = Pt(16); tf2.paragraphs[0].font.color.rgb = ACCENT_LIGHT; tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_transition(s, "fade", 2000)
    
    # ── Save ───────────────────────────────────────────────
    out = "/home/soang/.openclaw/workspace/smartvn-report/slides/OCR-Presentation.pptx"
    prs.save(out)
    print(f"✅ Saved: {out} ({len(prs.slides)} slides)")

if __name__ == "__main__":
    main()
