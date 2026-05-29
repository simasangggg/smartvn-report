#!/usr/bin/env python3
"""
Create SmartVN E-Commerce PowerPoint Presentation
Professional styled slides with images from PDF
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────
PRIMARY = RGBColor(0x1B, 0x2A, 0x4A)      # Deep navy
ACCENT = RGBColor(0x2E, 0x86, 0xDE)       # Vibrant blue
ACCENT_LIGHT = RGBColor(0xD6, 0xEA, 0xF8) # Light blue
SUCCESS = RGBColor(0x27, 0xAE, 0x60)      # Green
WARNING = RGBColor(0xF3, 0x9C, 0x12)      # Orange
DANGER = RGBColor(0xE7, 0x4C, 0x3C)       # Red
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
GRAY_TEXT = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)

IMG_DIR = "/home/soang/.openclaw/workspace/smartvn-report/images"

def add_bg(slide, color=PRIMARY):
    """Add solid background color to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color, left=0, top=0, width=None, height=None):
    """Add a colored rectangle shape as background element."""
    if width is None:
        width = Inches(13.333)
    if height is None:
        height = Inches(7.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_accent_bar(slide, top=0, height=Inches(0.08)):
    """Add accent color bar at top of slide."""
    return add_shape_bg(slide, ACCENT, top=top, height=height)

def add_footer(slide, text="SmartVN — Hệ thống Thương mại Điện tử Microservices"):
    """Add footer bar."""
    bar = add_shape_bg(slide, PRIMARY, top=Inches(7.1), height=Inches(0.4))
    tf = bar.text_frame
    tf.text = text
    for p in tf.paragraphs:
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

def add_title_slide(prs, title, subtitle="", date="Tháng 6, 2026"):
    """Create a title/cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_bg(slide, PRIMARY)
    
    # Accent circle decoration
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(5), Inches(5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x2E, 0x86, 0xDE)
    circle.fill.fore_color.brightness = 0.3
    circle.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = ACCENT_LIGHT
        p2.space_before = Pt(12)
    
    # Date
    p3 = tf.add_paragraph()
    p3.text = date
    p3.font.size = Pt(14)
    p3.font.color.rgb = GRAY_TEXT
    p3.space_before = Pt(24)
    
    # Institute
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    p4 = tf2.paragraphs[0]
    p4.text = "Học viện Công nghệ Bưu chính Viễn thông\nKhoa Công nghệ Thông tin 2"
    p4.font.size = Pt(14)
    p4.font.color.rgb = ACCENT_LIGHT
    
    return slide

def add_section_slide(prs, section_title, section_num=""):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, ACCENT)
    
    # Section number
    if section_num:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(2), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = section_num
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.color.brightness = -0.3
    
    # Section title
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(10), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = section_title
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    
    return slide

def add_content_slide(prs, title, bullets, image_path=None, image_position="right"):
    """Create a content slide with optional image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_footer(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # Content area
    if image_path and os.path.exists(image_path):
        if image_position == "right":
            # Text on left, image on right
            txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.5), Inches(5.5))
            img_left = Inches(7.2)
            img_top = Inches(1.5)
            img_width = Inches(5.5)
            img_height = Inches(5)
        else:
            # Image on top, text below
            txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(3))
            img_left = Inches(2)
            img_top = Inches(1.2)
            img_width = Inches(9)
            img_height = Inches(2.5)
        
        try:
            slide.shapes.add_picture(image_path, img_left, img_top, img_width, img_height)
        except Exception as e:
            print(f"Error adding image {image_path}: {e}")
    else:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))
    
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        
        if bullet.startswith("##"):
            # Sub-header
            p.text = bullet[2:].strip()
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.space_before = Pt(16)
        elif bullet.startswith("- "):
            # Bullet point
            p.text = "• " + bullet[2:]
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(6)
        else:
            p.text = bullet
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(6)
    
    return slide

def add_image_slide(prs, title, image_path, caption=""):
    """Create a full-image slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_footer(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # Image
    if os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.2))
        except Exception as e:
            print(f"Error adding image {image_path}: {e}")
    
    # Caption
    if caption:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = caption
        p2.font.size = Pt(12)
        p2.font.italic = True
        p2.font.color.rgb = GRAY_TEXT
        p2.alignment = PP_ALIGN.CENTER
    
    return slide

def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Create a two-column comparison slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_footer(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # Left column header
    left_box = add_shape_bg(slide, ACCENT, left=Inches(0.5), top=Inches(1.2), width=Inches(6), height=Inches(0.5))
    ltf = left_box.text_frame
    ltf.paragraphs[0].text = left_title
    ltf.paragraphs[0].font.size = Pt(16)
    ltf.paragraphs[0].font.bold = True
    ltf.paragraphs[0].font.color.rgb = WHITE
    ltf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Left column content
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6), Inches(4.8))
    ltf2 = left_content.text_frame
    ltf2.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = ltf2.paragraphs[0]
        else:
            p = ltf2.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(6)
    
    # Right column header
    right_box = add_shape_bg(slide, SUCCESS, left=Inches(6.8), top=Inches(1.2), width=Inches(6), height=Inches(0.5))
    rtf = right_box.text_frame
    rtf.paragraphs[0].text = right_title
    rtf.paragraphs[0].font.size = Pt(16)
    rtf.paragraphs[0].font.bold = True
    rtf.paragraphs[0].font.color.rgb = WHITE
    rtf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Right column content
    right_content = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(6), Inches(4.8))
    rtf2 = right_content.text_frame
    rtf2.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = rtf2.paragraphs[0]
        else:
            p = rtf2.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(6)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Create a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_footer(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # Table
    cols = len(headers)
    table_rows = len(rows) + 1
    table = slide.shapes.add_table(table_rows, cols, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5)).table
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT_LIGHT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_TEXT
    
    return slide

def create_presentation():
    """Main function to create the presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    img = lambda name: os.path.join(IMG_DIR, name)
    
    # ── Slide 1: Title ─────────────────────────────────────
    add_title_slide(prs,
        "HỆ THỐNG THƯƠNG MẠI ĐIỆN TỬ\nSMARTVN",
        "Kiến trúc Microservices với Spring Boot, Spring Cloud & Redis\nỨng dụng Mobile Android & Recommendation System",
        "Thành phố Hồ Chí Minh, Tháng 6 năm 2026"
    )
    
    # ── Slide 2: Outline ───────────────────────────────────
    add_content_slide(prs, "NỘI DUNG TRÌNH BÀY", [
        "## 1. Tổng quan đề tài",
        "## 2. Kiến trúc hệ thống",
        "## 3. Công nghệ sử dụng",
        "## 4. Thiết kế cơ sở dữ liệu",
        "## 5. Recommendation System",
        "## 6. Mobile Application",
        "## 7. Kiểm thử & Đánh giá",
        "## 8. Kết luận",
    ])
    
    # ── Section 1: Tổng quan ───────────────────────────────
    add_section_slide(prs, "TỔNG QUAN ĐỀ TÀI", "01")
    
    add_content_slide(prs, "BỐI CẢNH & BÀI TOÁN", [
        "## Thương mại điện tử tại Việt Nam",
        "- Thị trường TMĐT toàn cầu > 6.000 tỷ USD (2025)",
        "- Việt Nam: Shopee, Lazada, Tiki, Sendo cạnh tranh mạnh",
        "",
        "## Thách thức đặt ra",
        "- Mở rộng linh hoạt cho lượng lớn người dùng",
        "- Đảm bảo tính sẵn sàng cao (fault tolerance)",
        "- Bảo mật thông tin & giao dịch thanh toán",
        "- Duy trì và phát triển từng phần độc lập",
        "",
        "## Hướng tiếp cận",
        "- Kiến trúc Microservices với Spring Boot 3.3 + Spring Cloud",
        "- Mobile Android native (Java, MVVM)",
        "- AI Recommendation (PhoBERT + ALS)",
    ])
    
    add_content_slide(prs, "MỤC TIÊU ĐỀ TÀI", [
        "- Thiết kế hệ thống TMĐT SmartVN với Microservices",
        "- Phát triển ứng dụng Mobile Android native (MVVM)",
        "- Triển khai 7 microservices: Eureka, Config, Gateway, User, Product, Order, Admin",
        "- Cache Redis cho Product Service (giảm latency 5.5x)",
        "- Xác thực JWT + OAuth2 (Google, GitHub)",
        "- Thanh toán trực tuyến VNPay",
        "- Circuit Breaker (Resilience4j) cho fault tolerance",
        "- AI Recommendation: PhoBERT + ALS Hybrid",
        "- Docker Compose deployment",
        "- Benchmark hiệu năng với Grafana k6",
    ], image_path=img("img-018.jpg"), image_position="right")
    
    add_image_slide(prs, "QUY TRÌNH MUA HÀNG — TRƯỚC KHI CÓ HỆ THỐNG",
        img("img-014.jpg"),
        "Kịch bản A: Quy trình thủ công — giấy tờ, Excel, gọi điện, nhập tay")
    
    add_image_slide(prs, "QUY TRÌNH MUA HÀNG — SAU KHI CÓ HỆ THỐNG",
        img("img-015.jpg"),
        "Kịch bản B: Quy trình tự động hoá qua Web App — đặt hàng, thanh toán, theo dõi online")
    
    # ── Section 2: Kiến trúc ───────────────────────────────
    add_section_slide(prs, "KIẾN TRÚC HỆ THỐNG", "02")
    
    add_image_slide(prs, "SƠ ĐỒ KIẾN TRÚC TỔNG THỂ — CONTAINER VIEW",
        img("img-018.jpg"),
        "Kiến trúc Microservices: API Gateway → Service Discovery → Business Services → Data Layer")
    
    add_two_column_slide(prs, "SO SÁNH: MONOLITHIC vs MICROSERVICES",
        "Monolithic", [
            "Một khối code duy nhất",
            "Khó mở rộng theo chiều ngang",
            "Deploy toàn bộ khi thay đổi",
            "Một lỗi nhỏ → sập toàn hệ thống",
            "Công nghệ đơn nhất",
            "Dễ phát triển ban đầu",
        ],
        "Microservices", [
            "Nhiều dịch vụ nhỏ, độc lập",
            "Mở rộng linh hoạt từng service",
            "Deploy từng phần độc lập",
            "Fault isolation — lỗi service này không ảnh hưởng khác",
            "Polyglot — mỗi service chọn công nghệ phù hợp",
            "Phức tạp hơn khi quản lý",
        ])
    
    add_content_slide(prs, "CÁC MICROSERVICES TRONG HỆ THỐNG", [
        "## Infrastructure Services",
        "- Eureka Server (port 8761) — Service Discovery & Registry",
        "- Config Server (port 8888) — Quản lý cấu hình tập trung",
        "- API Gateway (port 8080) — Routing, JWT validation, CORS",
        "",
        "## Business Services",
        "- User Service (port 8081) — Quản lý người dùng, xác thực JWT/OAuth2",
        "- Product Service (port 8082) — Quản lý sản phẩm, cache Redis",
        "- Order Service (port 8083) — Quản lý đơn hàng, thanh toán VNPay",
        "- Admin Service (port 8084) — Quản trị, dashboard",
        "",
        "## Data Layer",
        "- MySQL 8.0 (port 3306) — Cơ sở dữ liệu chính",
        "- Redis 7 (port 6379) — Bộ nhớ đệm",
    ])
    
    # ── Section 3: Công nghệ ───────────────────────────────
    add_section_slide(prs, "CÔNG NGHỆ SỬ DỤNG", "03")
    
    add_table_slide(prs, "BẢNG TỔNG HỢP CÔNG NGHỆ",
        ["Layer", "Công nghệ", "Vai trò"],
        [
            ["Backend", "Spring Boot 3.3 + Spring Cloud", "Microservices framework"],
            ["Service Discovery", "Netflix Eureka", "Service registry & discovery"],
            ["API Gateway", "Spring Cloud Gateway", "Routing, JWT, CORS"],
            ["Database", "MySQL 8.0", "RDBMS chính"],
            ["Cache", "Redis 7", "Bộ nhớ đệm, giảm latency"],
            ["Security", "Spring Security + JWT + OAuth2", "Xác thực & phân quyền"],
            ["Payment", "VNPay Gateway", "Thanh toán trực tuyến"],
            ["Resilience", "Resilience4j Circuit Breaker", "Fault tolerance"],
            ["Mobile", "Android Java + MVVM", "Ứng dụng di động"],
            ["Frontend", "React + Vite", "Web UI (Customer & Admin)"],
            ["AI/ML", "PhoBERT + ALS", "Product recommendation"],
            ["Container", "Docker Compose", "Deployment & orchestration"],
            ["Testing", "Grafana k6", "Performance benchmark"],
        ])
    
    add_content_slide(prs, "SPRING CLOUD ECOSYSTEM", [
        "## Service Discovery — Netflix Eureka",
        "- Service tự đăng ký khi khởi động",
        "- Client lookup để tìm service khác",
        "- Health check tự động, loại bỏ service chết",
        "",
        "## API Gateway — Spring Cloud Gateway",
        "- Điểm vào duy nhất cho tất cả request",
        "- Route dựa trên path: /api/users/** → User Service",
        "- JWT validation trước khi chuyển tiếp",
        "- CORS handling, rate limiting",
        "",
        "## Circuit Breaker — Resilience4j",
        "- Open circuit khi service downstream fail",
        "- Fallback response thay vì timeout treo",
        "- Half-open để thử service đã phục hồi chưa",
    ])
    
    add_content_slide(prs, "REDIS CACHING — HIỆU NĂNG VƯỢT TRỘI", [
        "## Cache-aside Pattern",
        "- Ứng dụng kiểm tra Redis trước khi query MySQL",
        "- Cache miss → query MySQL → set Redis → trả về",
        "- Cache hit → trả về trực tiếp (không cần query DB)",
        "",
        "## Kết quả Benchmark (k6)",
        "- Không cache: 47.81ms avg response time",
        "- Có Redis cache: 8.66ms avg (giảm 5.5 lần!)",
        "- Throughput: 196.8 → 248.5 req/s (tăng 26%)",
        "- Error rate: giảm xuống 0%",
        "",
        "## Chiến lược TTL",
        "- Product list: TTL 5 phút",
        "- Product detail: TTL 10 phút",
        "- Invalidate khi có update (write-through)",
    ])
    
    add_content_slide(prs, "XÁC THỰC ĐA TẦNG — JWT & OAuth2", [
        "## JWT (JSON Web Token)",
        "- Access Token: hết hạn 15 phút, lưu memory",
        "- Refresh Token: hết hạn 7 ngày, lưu HTTP-only cookie",
        "- Interceptor tự động refresh khi access token hết hạn",
        "",
        "## OAuth2 (Google, GitHub)",
        "- Frontend redirect đến trang đăng nhập Google",
        "- Google callback với authorization code",
        "- Backend đổi code lấy user info, tạo/tìm tài khoản",
        "- Trả JWT token cho frontend",
        "",
        "## API Key cho Inter-service Communication",
        "- Mỗi service có API key riêng",
        "- Header X-API-Key khi gọi nội bộ",
        "- Không cần JWT cho service-to-service",
    ])
    
    # ── Section 4: Database ────────────────────────────────
    add_section_slide(prs, "THIẾT KẾ CƠ SỞ DỮ LIỆU", "04")
    
    add_image_slide(prs, "SƠ ĐỒ ERD — CƠ SỞ DỮ LIỆU",
        img("img-070.jpg"),
        "Entity Relationship Diagram: users, roles, addresses và các bảng liên quan")
    
    add_table_slide(prs, "CÁC BẢNG CHÍNH TRONG HỆ THỐNG",
        ["Bảng", "Service", "Mô tả"],
        [
            ["users", "User Service", "Thông tin người dùng, xác thực"],
            ["roles", "User Service", "Phân quyền: Admin, Staff, Customer"],
            ["addresses", "User Service", "Địa chỉ giao hàng"],
            ["products", "Product Service", "Danh mục sản phẩm điện tử"],
            ["product_images", "Product Service", "Hình ảnh sản phẩm"],
            ["categories", "Product Service", "Phân loại sản phẩm"],
            ["orders", "Order Service", "Đơn hàng và trạng thái"],
            ["order_items", "Order Service", "Chi tiết từng mặt hàng"],
            ["payments", "Order Service", "Giao dịch VNPay"],
            ["reviews", "Admin Service", "Đánh giá sản phẩm"],
        ])
    
    # ── Section 5: Recommendation ──────────────────────────
    add_section_slide(prs, "RECOMMENDATION SYSTEM", "05")
    
    add_content_slide(prs, "HỆ THỐNG GỢI Ý SẢN PHẨM — HYBRID APPROACH", [
        "## Tại sao cần Recommendation?",
        "- Người dùng không muốn duyệt hàng nghìn sản phẩm",
        "- Cá nhân hoá trải nghiệm mua sắm",
        "- Tăng doanh thu qua cross-sell & up-sell",
        "",
        "## Hybrid = Content-Based + Collaborative Filtering",
        "- Content-Based: PhoBERT phân tích mô tả sản phẩm",
        "- Collaborative: ALS học từ lịch sử mua hàng",
        "- Kết hợp điểm từ cả 2 model → gợi ý tốt nhất",
        "",
        "## Xử lý Cold Start",
        "- User mới → dùng Content-Based (không cần lịch sử)",
        "- Sản phẩm mới → dùng Collaborative (từ user tương tự)",
    ], image_path=img("img-019.jpg"), image_position="right")
    
    add_image_slide(prs, "PHOBERT — CONTENT-BASED RECOMMENDATION",
        img("img-004.jpg"),
        "PhoBERT (VinAI): Pre-trained Vietnamese language model cho phân tích nội dung sản phẩm")
    
    add_image_slide(prs, "ALS — COLLABORATIVE FILTERING PIPELINE",
        img("img-030.jpg"),
        "ALS Training Pipeline: Data extraction → Transformation → Model training → Storage")
    
    add_image_slide(prs, "HYBRID RECOMMENDATION — SEQUENCE DIAGRAM",
        img("img-040.jpg"),
        "Luồng gợi ý: Frontend → API Gateway → Recommend Service → PhoBERT + ALS → Hybrid Merge → Response")
    
    # ── Section 6: Mobile ──────────────────────────────────
    add_section_slide(prs, "ỨNG DỤNG MOBILE ANDROID", "06")
    
    add_content_slide(prs, "KIẾN TRÚC ỨNG DỤNG MOBILE — MVVM", [
        "## Tech Stack",
        "- Java + Android SDK",
        "- Retrofit 2 cho REST API calls",
        "- Room Database cho offline-first",
        "- LiveData + ViewModel cho reactive UI",
        "- Material Design 3 cho giao diện hiện đại",
        "",
        "## Features",
        "- Đăng ký / Đăng nhập (JWT + OAuth2 Google)",
        "- Duyệt sản phẩm, tìm kiếm, lọc theo danh mục",
        "- Giỏ hàng (offline + sync với server)",
        "- Đặt hàng & thanh toán VNPay",
        "- Theo dõi trạng thái đơn hàng",
        "- Quản lý tài khoản & địa chỉ",
    ], image_path=img("img-019.jpg"), image_position="right")
    
    add_content_slide(prs, "USE CASE DIAGRAM — ỨNG DỤNG MOBILE", [
        "## Khách vãng lai (Guest)",
        "- Xem danh sách sản phẩm",
        "- Tìm kiếm sản phẩm",
        "- Xem chi tiết sản phẩm",
        "- Xem đánh giá",
        "",
        "## Khách hàng (Logged-in)",
        "- Tất cả quyền của Guest +",
        "- Nhận gợi ý sản phẩm (AI-Powered)",
        "- Quản lý giỏ hàng",
        "- Đặt hàng & thanh toán",
        "- Đánh giá sản phẩm",
        "- Quản lý tài khoản & địa chỉ",
        "- Xem lịch sử đơn hàng",
    ], image_path=img("img-019.jpg"), image_position="right")
    
    # ── Section 7: UI Screens ──────────────────────────────
    add_section_slide(prs, "GIAO DIỆN HỆ THỐNG", "07")
    
    add_image_slide(prs, "GIAO DIỆN ĐĂNG NHẬP",
        img("img-080.jpg"),
        "Login page: Email/Password + Social login (Google, GitHub)")
    
    add_image_slide(prs, "GIAO DIỆN DANH SÁCH SẢN PHẨM",
        img("img-085.jpg"),
        "Product listing: Category filter, price range, product cards với Sale badge")
    
    add_image_slide(prs, "GIAO DIỆN QUẢN LÝ TÀI KHOẢN",
        img("img-090.jpg"),
        "User profile: Profile Info, Addresses, Orders, My Reviews tabs")
    
    add_image_slide(prs, "GIAO DIỆN ĐẶT HÀNG THÀNH CÔNG",
        img("img-095.jpg"),
        "Order confirmation: Order ID, Total Amount, Payment Method (VNPay), Shipping address")
    
    add_image_slide(prs, "ADMIN DASHBOARD — QUẢN LÝ NGƯỜI DÙNG",
        img("img-100.jpg"),
        "Admin panel: User management với Role, Status, Warnings, Actions")
    
    add_image_slide(prs, "ADMIN DASHBOARD — CHỈNH SỬA SẢN PHẨM",
        img("img-105.jpg"),
        "Product edit: Name, Brand, Color, Description, Technical specs (weight, RAM, battery...)")
    
    add_image_slide(prs, "ADMIN DASHBOARD — QUẢN LÝ ĐÁNH GIÁ",
        img("img-109.jpg"),
        "Review management: Star ratings, review text, user info, delete action")
    
    # ── Section 8: Testing ─────────────────────────────────
    add_section_slide(prs, "KIỂM THỬ & ĐÁNH GIÁ", "08")
    
    add_table_slide(prs, "KẾT QUẢ BENCHMARK — REDIS CACHE IMPACT",
        ["Metric", "Không Cache", "Redis Cache", "Cải thiện"],
        [
            ["Avg Response Time", "47.81ms", "8.66ms", "5.5x nhanh hơn"],
            ["Throughput", "196.8 req/s", "248.5 req/s", "+26%"],
            ["Error Rate", "> 0%", "0%", "Hoàn hảo"],
            ["P95 Latency", "~120ms", "~25ms", "4.8x nhanh hơn"],
        ])
    
    add_content_slide(prs, "ĐÁNH GIÁ TỔNG QUAN", [
        "## Ưu điểm hệ thống",
        "- Microservices cho phép scale và deploy độc lập",
        "- Redis cache giảm latency 5.5 lần",
        "- JWT + OAuth2 bảo mật mạnh",
        "- Circuit Breaker đảm bảo fault tolerance",
        "- AI Recommendation cá nhân hoá trải nghiệm",
        "- Docker Compose deployment đơn giản",
        "",
        "## Hạn chế",
        "- Chưa triển khai production thực tế",
        "- Chưa có CI/CD pipeline",
        "- Push notification (FCM) chưa hoàn thiện",
        "- Chưa có monitoring (Prometheus/Grafana)",
        "- Frontend chưa tối ưu SEO",
    ])
    
    # ── Section 9: Kết luận ────────────────────────────────
    add_section_slide(prs, "KẾT LUẬN & HƯỚNG PHÁT TRIỂN", "09")
    
    add_content_slide(prs, "KẾT QUẢ ĐẠT ĐƯỢC", [
        "- Hệ thống TMĐT SmartVN hoàn chỉnh với 7 microservices",
        "- Ứng dụng Mobile Android native đầy đủ tính năng",
        "- AI Recommendation System (PhoBERT + ALS Hybrid)",
        "- Hiệu năng ấn tượng: 5.5x nhanh hơn với Redis cache",
        "- Bảo mật đa tầng: JWT + OAuth2 + API Key",
        "- Fault tolerance với Circuit Breaker",
        "- Docker Compose deployment一键启动",
    ])
    
    add_content_slide(prs, "HƯỚNG PHÁT TRIỂN", [
        "## Ngắn hạn",
        "- Triển khai CI/CD pipeline (GitHub Actions)",
        "- Thêm monitoring (Prometheus + Grafana)",
        "- Push notification qua Firebase Cloud Messaging",
        "- Tối ưu SEO cho Customer Frontend",
        "",
        "## Dài hạn",
        "- Migrate sang Kubernetes cho orchestration",
        "- Thêm Elasticsearch cho full-text search",
        "- Tích hợp chatbot AI hỗ trợ khách hàng",
        "- Multi-language support (i18n)",
        "- Mobile app iOS (Flutter cross-platform)",
    ])
    
    # ── Slide cuối: Cảm ơn ─────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY)
    
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-2), Inches(7), Inches(7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.fill.fore_color.brightness = 0.3
    circle.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "CẢM ƠN THẦY VÀ CÁC BẠN!"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Xin cảm ơn!"
    p2.font.size = Pt(24)
    p2.font.color.rgb = ACCENT_LIGHT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(10), Inches(1))
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = "Học viện Công nghệ Bưu chính Viễn thông\nKhoa Công nghệ Thông tin 2"
    p3.font.size = Pt(16)
    p3.font.color.rgb = ACCENT_LIGHT
    p3.alignment = PP_ALIGN.CENTER
    
    # ── Save ───────────────────────────────────────────────
    output_path = "/home/soang/.openclaw/workspace/smartvn-report/SmartVN-Slides.pptx"
    prs.save(output_path)
    print(f"✅ Presentation saved: {output_path}")
    print(f"   Slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
